import pypdf
from pathlib import Path
import re


def _chunk_from_pages(pages, source_name, chunk_size=1000, overlap=200):
    """Core chunking: flatten (page_num, text) pairs then split into overlapping sentence chunks."""
    full_text = ""
    char_to_page = []
    for page_num, text in pages:
        full_text += text
        char_to_page.extend([page_num] * len(text))

    if not full_text.strip():
        return []

    # lookbehind on .!? keeps punctuation attached to the sentence that ends with it
    sentences = re.split(r'(?<=[.!?])\s+', full_text)

    # record the start character of each sentence so we can map it back to a page number
    sentence_positions = []
    pos = 0
    for sentence in sentences:
        sentence_positions.append(pos)
        pos += len(sentence) + 1  # +1 for the space that was split on

    chunks = []
    i = 0
    while i < len(sentences):
        chunk_sentences = []
        chunk_len = 0

        j = i
        while j < len(sentences) and chunk_len + len(sentences[j]) <= chunk_size:
            chunk_sentences.append(sentences[j])
            chunk_len += len(sentences[j]) + 1
            j += 1

        # if a single sentence exceeds chunk_size, include it anyway
        if not chunk_sentences:
            chunk_sentences.append(sentences[i])
            j = i + 1

        chunk_text = " ".join(chunk_sentences)
        char_pos = sentence_positions[i]
        # clamp to last index in case rounding puts us one past the end
        page = char_to_page[min(char_pos, len(char_to_page) - 1)]

        chunks.append({
            "text": chunk_text,
            "source": source_name,
            "page": page,
            "chunk_index": len(chunks),
        })

        # find overlap: step back from j until we've covered ~overlap chars
        overlap_len = 0
        step_back = j - 1
        while step_back > i and overlap_len < overlap:
            overlap_len += len(sentences[step_back]) + 1
            step_back -= 1
        i = step_back + 1

    return chunks


# ── Extractors ────────────────────────────────────────────────────────────────

def _extract_pdf(file_path):
    """Extract text page-by-page from a PDF."""
    pdf = pypdf.PdfReader(file_path)
    return [(n + 1, page.extract_text() or "") for n, page in enumerate(pdf.pages)]


def _extract_plain_text(file_path):
    """Read a plain text file as a single page."""
    text = Path(file_path).read_text(encoding="utf-8", errors="replace")
    return [(1, text)]


def _extract_epub(file_path):
    """Extract text chapter-by-chapter from an EPUB ebook."""
    import ebooklib
    from ebooklib import epub
    from html.parser import HTMLParser

    class _Strip(HTMLParser):
        def __init__(self):
            super().__init__()
            self._buf = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip = True

        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip and data.strip():
                self._buf.append(data.strip())

    book = epub.read_epub(str(file_path))
    pages = []
    for chapter_num, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT), start=1):
        content = item.get_content().decode("utf-8", errors="replace")
        parser = _Strip()
        parser.feed(content)
        text = " ".join(parser._buf)
        if text.strip():
            pages.append((chapter_num, text))
    return pages or [(1, "")]


def _extract_pptx(file_path):
    """Extract text slide-by-slide from a PowerPoint presentation (.pptx)."""
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = [
            para.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        ]
        if texts:
            pages.append((slide_num, "\n".join(texts)))
    return pages or [(1, "")]


def _extract_ppt(file_path):
    raise ValueError(
        "Legacy .ppt binary format is not supported. "
        "Open the file in PowerPoint and save it as .pptx first."
    )


def _extract_xlsx(file_path):
    """Extract each sheet from an .xlsx workbook as a separate page."""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    pages = []
    for page_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows = [
            "\t".join("" if c is None else str(c) for c in row)
            for row in ws.iter_rows(values_only=True)
        ]
        text = "\n".join(r for r in rows if r.strip())
        if text:
            pages.append((page_num, text))
    wb.close()
    return pages or [(1, "")]


def _extract_xls(file_path):
    """Extract each sheet from a legacy .xls workbook as a separate page."""
    import xlrd
    wb = xlrd.open_workbook(file_path)
    pages = []
    for page_num, sheet in enumerate(wb.sheets(), start=1):
        rows = [
            "\t".join(str(c) for c in sheet.row_values(r))
            for r in range(sheet.nrows)
        ]
        text = "\n".join(r for r in rows if r.strip())
        if text:
            pages.append((page_num, text))
    return pages or [(1, "")]


def _extract_docx(file_path):
    """Extract paragraph text from a Word document (.docx), grouped into sections."""
    import docx
    doc = docx.Document(file_path)
    # Group every 20 paragraphs as one "page" so long documents stay navigable
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    for page_num, start in enumerate(range(0, max(len(paras), 1), 20), start=1):
        text = "\n".join(paras[start:start + 20])
        if text:
            pages.append((page_num, text))
    return pages or [(1, "")]


def _extract_html(file_path):
    """Strip HTML tags and extract visible text using the built-in html.parser."""
    from html.parser import HTMLParser

    class _Extractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self._buf = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style", "head"):
                self._skip = True

        def handle_endtag(self, tag):
            if tag in ("script", "style", "head"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip:
                stripped = data.strip()
                if stripped:
                    self._buf.append(stripped)

    content = Path(file_path).read_text(encoding="utf-8", errors="replace")
    parser = _Extractor()
    parser.feed(content)
    return [(1, " ".join(parser._buf))]


def _extract_csv(file_path):
    """Parse a CSV into readable 'key: value' rows so embeddings capture column context."""
    import csv
    rows = []
    with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        headers = None
        for row in reader:
            if headers is None:
                headers = row
                rows.append(", ".join(row))
            else:
                rows.append(", ".join(f"{h}: {v}" for h, v in zip(headers, row)))
    return [(1, "\n".join(rows))]


def _extract_structured_text(file_path):
    """Read structured text formats (JSON, YAML, TOML, XML) as a single page of raw text."""
    text = Path(file_path).read_text(encoding="utf-8", errors="replace")
    return [(1, text)]


_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".txt":  _extract_plain_text,
    ".md":   _extract_plain_text,
    ".log":  _extract_plain_text,
    ".rst":  _extract_plain_text,
    ".ini":  _extract_plain_text,
    ".cfg":  _extract_plain_text,
    ".conf": _extract_plain_text,
    ".toml": _extract_structured_text,
    ".yaml": _extract_structured_text,
    ".yml":  _extract_structured_text,
    ".json": _extract_structured_text,
    ".xml":  _extract_structured_text,
    ".csv":  _extract_csv,
    ".html": _extract_html,
    ".htm":  _extract_html,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xls,
    ".pptx": _extract_pptx,
    ".ppt":  _extract_ppt,
    ".epub": _extract_epub,
}

# Exported so other modules can validate without importing private internals
SUPPORTED_EXTENSIONS = frozenset(_EXTRACTORS)


def chunk_file(file_path, chunk_size=1000, overlap=200):
    """Chunk any supported file into overlapping sentence-aware chunks with source and page metadata."""
    ext = Path(file_path).suffix.lower()
    if ext not in _EXTRACTORS:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    pages = _EXTRACTORS[ext](file_path)
    return _chunk_from_pages(pages, Path(file_path).name, chunk_size, overlap)


# kept for backward compatibility
def chunk_pdf(file_path, chunk_size=1000, overlap=200):
    """Split a PDF into overlapping sentence-aware chunks with source and page metadata."""
    return chunk_file(file_path, chunk_size, overlap)


if __name__ == "__main__":
    file_path = Path(__file__).parent.parent.parent / "resources" / "Attention-Is-All-You-Need.pdf"
    chunks = chunk_file(file_path)
    for idx, chunk in enumerate(chunks):
        print(f"Chunk {idx+1}:\n{chunk}\n")
